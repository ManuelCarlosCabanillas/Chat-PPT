import streamlit as st
from PyPDF2 import PdfReader
import openai
import re
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings, HuggingFaceInstructEmbeddings
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from pptx import Presentation



def main():
    st.title('🧠 Chat PPT')
    st.markdown('Choose which version of the application to run:')
    st.markdown('**- Context**: for asking directly gpt 3.5 TURBO model - precission is higher but it is limited to 16K Tokens')
    st.markdown('**- Embeddings**: can load several pdfs documments and transform them into vectors but precission is a little bit lower')

    # Choose version
    version = st.selectbox("Choose version", ["Select","Context", "Embeddings"])
    
    if version == "Context":
        context()
        # Initial version
        pass
    elif version == "Embeddings":
        # Alternative version
        embeddings()
    else:
        st.warning("Please select a version of the application to run.")

def context():
    openai.api_key = st.secrets["openai_api_key"]    
    
    # This is a helper function to read PPTs
    def read_ppt(ppt, slides):
        text = ""
        for slide in slides:
            for shape in ppt.slides[slide-1].shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
                            text += "\n"
        return text    
    def ask_gpt3(question, context, temperature, max_tokens, top_p, frequency_penalty, role):
        message = [
            {"role": "system", "content": "You have the following information from the ppt: "+context},
            {"role": role, "content": question}
        ]
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo-16k",
            messages=message,
            temperature=temperature,
            max_tokens=max_tokens,
            top_p=top_p,
            frequency_penalty=frequency_penalty
        )
        return response['choices'][0]['message']['content']
    
    st.title('Ask directly GPT on the context given by a ppt document')

    # Configure the file uploader
# Configure the file uploader
    uploaded_file = st.file_uploader("Upload your PPT file", type="pptx")
    
    if uploaded_file is not None:
        # Read the slides of the PPT
        ppt = Presentation(uploaded_file)
        num_slides = len(ppt.slides)
    
        # Configure the parameters of the OpenAI API
        st.sidebar.title('🛠️ OpenAI API Configuration')
        slide_selection = st.sidebar.multiselect('Slides', options=range(1, num_slides+1), default=range(1, num_slides+1))        
        st.sidebar.markdown("<small>Select the slides you want to use as context</small>", unsafe_allow_html=True)
        temperature = st.sidebar.slider('Temperature', min_value=0.0, max_value=1.0, value=0.5)
        st.sidebar.markdown("""<small>Temperature determines the randomness of the AI's responses. A higher value will make the responses more diverse, but also riskier.</small>""", unsafe_allow_html=True)
        max_tokens = st.sidebar.slider('Max Tokens', min_value=10, max_value=500, value=500)
        st.sidebar.markdown("""<small>Max tokens limit the length of the AI's response.</small>""", unsafe_allow_html=True)
        top_p = st.sidebar.slider('Top P', min_value=0.0, max_value=1.0, value=0.9)
        st.sidebar.markdown("""<small>Top P is the cumulative probability by the highest-ranking words, which affects the diversity of the response.</small>""", unsafe_allow_html=True)
        frequency_penalty = st.sidebar.slider('Frequency Penalty', min_value=-2.0, max_value=2.0, value=0.0)
        st.sidebar.markdown("""<small>Frequency penalty reduces the likelihood of frequent words.</small>""", unsafe_allow_html=True)
        role = st.sidebar.selectbox('Role', ('system', 'user', 'assistant'), index=2)
        st.sidebar.markdown("""<small>The role defines the behavior of the chatbot.</small>""", unsafe_allow_html=True)
    
        # Convert the slide selection to 0-indexed
        slide_selection = [slide-1 for slide in slide_selection]
    
        # Use the helper function to extract the text from the selected slides
        context = read_ppt(ppt, slide_selection)
    
        # Create a text input field for the question
        question = st.text_input("Enter your question here")
    
        if st.button('Ask the question'):
            if question:
                try:
                    # Use the OpenAI API to get an answer
                    response = ask_gpt3(question, context, temperature, max_tokens, top_p, frequency_penalty, role)
    
                    # Display the answer
                    st.markdown("**Answer:**")
                    st.markdown(response)
                except openai.error.InvalidRequestError as e:
                    # Extract the number of requested tokens and the maximum allowed from the error message
                    max_tokens, tokens_requested = re.findall(r'\d+', str(e))
                    st.error(f"You have requested {tokens_requested} tokens when the maximum allowed is {max_tokens}. Please reduce the number of slides in the configuration bar.")
            else:
                st.warning('Please enter a question.')


def embeddings():
    openai.api_key = st.secrets["openai_api_key"]
    def get_content_from_files(docs):
        text = ""
        for doc in docs:
            if doc.name.endswith(".pdf"):
                pdf_reader = PdfReader(doc)
                for page in pdf_reader.pages:
                    text += page.extract_text()
            elif doc.name.endswith(".pptx"):
                ppt = Presentation(doc)
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame"):
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    text += run.text
                                    text += "\n"
        return text

    def split_text_into_chunks(text):
        text_splitter = CharacterTextSplitter(
            separator="\n",
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len
        )
        chunks = text_splitter.split_text(text)
        return chunks

    def generate_vectorstore(text_chunks):
        embeddings = OpenAIEmbeddings(openai_api_key=openai.api_key)
        vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
        return vectorstore

    def create_conversation_chain(vectorstore):
        llm = ChatOpenAI(openai_api_key=openai.api_key)
        memory = ConversationBufferMemory(
            memory_key='chat_history', return_messages=True)
        conversation_chain = ConversationalRetrievalChain.from_llm(
            llm=llm,
            retriever=vectorstore.as_retriever(),
            memory=memory
        )
        return conversation_chain

    def handle_userinput(user_question):
        if st.session_state.conversation is not None:
            response = st.session_state.conversation({'question': user_question})
            st.session_state.chat_history = response['chat_history']

            for i, message in enumerate(st.session_state.chat_history):
                if i % 2 == 0:  # User's messages
                    st.write(message.content)
                else:  # Bot's messages
                    st.markdown(f'**{message.content}**')
        else:
            st.warning('Please process your documents before asking a question.')
    
    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None

    st.sidebar.header("Model Parameters")
    chunk_size = st.sidebar.slider("Chunk Size", 500, 2000, 1000)
    chunk_overlap = st.sidebar.slider("Chunk Overlap", 0, 500, 200)
    
    st.header("Upload and Process Documents")
    # Configure the file uploader
    pdf_docs = st.file_uploader(
    "Upload your PDFs and PPTs here and click on 'Process'", accept_multiple_files=True)

    if st.button("Process"):
        with st.spinner("Processing"):
            # get text from uploaded files
            raw_text = get_content_from_files(pdf_docs)

            # get the text chunks
            text_chunks = split_text_into_chunks(raw_text)

            # create vector store
            vectorstore = generate_vectorstore(text_chunks)

            # create conversation chain
            st.session_state.conversation = create_conversation_chain(vectorstore)
        
    if st.session_state.conversation is not None:
        st.header("Ask something about your documents")
        user_question = st.text_input("")
        if user_question:
            handle_userinput(user_question)

if __name__ == '__main__':
    main()
